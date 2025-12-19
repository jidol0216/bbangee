#!/usr/bin/env python3
"""
Face Tracking Presentation Generator
Markdown → PowerPoint 변환
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 색상 정의
BLUE = RGBColor(0, 112, 192)
DARK_BLUE = RGBColor(0, 51, 102)
ORANGE = RGBColor(255, 153, 0)
GREEN = RGBColor(0, 176, 80)
RED = RGBColor(255, 0, 0)
GRAY = RGBColor(128, 128, 128)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)


def add_title_slide(prs, title, subtitle=""):
    """타이틀 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경색 (진한 파랑)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 서브타이틀
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs, title, emoji="📌"):
    """섹션 구분 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BLUE
    background.line.fill.background()
    
    # 이모지
    emoji_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = emoji_box.text_frame
    p = tf.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(72)
    p.alignment = PP_ALIGN.CENTER
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, content_lines, bullet=True):
    """내용 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 타이틀 배경
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_BLUE
    title_bg.line.fill.background()
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if bullet and line.strip():
            p.text = "• " + line
        else:
            p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.space_after = Pt(12)
    
    return slide


def add_table_slide(prs, title, headers, rows):
    """테이블 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 타이틀 배경
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_BLUE
    title_bg.line.fill.background()
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 테이블
    cols = len(headers)
    table_rows = len(rows) + 1
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.5 * table_rows)
    
    table = slide.shapes.add_table(table_rows, cols, left, top, width, height).table
    
    # 헤더
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # 데이터
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_diagram_slide(prs, title, diagram_text):
    """다이어그램 슬라이드 (텍스트 기반)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 타이틀 배경
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_BLUE
    title_bg.line.fill.background()
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 다이어그램 (고정폭 폰트)
    diagram_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.2))
    tf = diagram_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = diagram_text
    p.font.size = Pt(12)
    p.font.name = "Consolas"
    p.font.color.rgb = BLACK
    
    return slide


def create_presentation():
    """전체 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ==================== 슬라이드 1: 타이틀 ====================
    add_title_slide(
        prs,
        "🤖 실시간 얼굴 추적 로봇 시스템",
        "Doosan M0609 + RealSense D435i + ROS2\n\nRokey Bootcamp Col2 Team"
    )
    
    # ==================== 슬라이드 2: 목차 ====================
    add_content_slide(prs, "📋 발표 목차", [
        "전체 시나리오 - 군 위병소 자동화 시스템",
        "프로젝트 소개 - 무엇을 만들었나?",
        "문제 정의 - 왜 어려운가?",
        "개발 히스토리 - 어떻게 해결했나? (Day 1~4)",
        "MPC → Joint 전환 - 왜 제어 방식을 바꿨나?",
        "시스템 아키텍처 - 전체 구조",
        "기술 선택 비교 - 왜 이 기술을 선택했나?",
        "핵심 알고리즘 - EKF, Joint Control",
        "결과 및 성과 - 무엇을 달성했나?",
        "미완성 모듈 & 향후 계획"
    ])
    
    # ==================== 슬라이드 3: 전체 시나리오 ====================
    add_section_slide(prs, "전체 시나리오", "🎯")
    
    add_diagram_slide(prs, "🏛️ 군 위병소 자동화 시스템 (ver1)", """
    장소: 군 부대 위병소  |  상황: 로봇팔이 초병 역할

    ┌─────────────────────────────────────────────────────────────────┐
    │ STEP 1: 권총 파지                                                │
    │   → 아무렇게 올려진 권총을 SAM3 기반 그리퍼로 파지                   │
    └─────────────────────────────────────────────────────────────────┘
                                    ▼
    ┌─────────────────────────────────────────────────────────────────┐
    │ STEP 2: 사격 위치 이동                                           │
    │   → 사로(射路)에 고정 → 다시 파지                                  │
    └─────────────────────────────────────────────────────────────────┘
                                    ▼
    ┌─────────────────────────────────────────────────────────────────┐
    │ STEP 3: 사주 경계 (★ 현재 구현된 부분)                            │
    │   → 헤드샷 트래킹 (YOLOv8-face + EKF + Joint Control)            │
    │   → 레이저 포인터 활성화 (가장 가까운 사람)                         │
    └─────────────────────────────────────────────────────────────────┘
                                    ▼
    ┌─────────────────────────────────────────────────────────────────┐
    │ STEP 4: 피아 식별                                                │
    │   → 아군/적군 클래스 구분 (학습 데이터 기반)                         │
    └─────────────────────────────────────────────────────────────────┘
    """)
    
    add_diagram_slide(prs, "🔀 피아식별 분기 로직", """
    ┌─────────────────────────────────────────────────────────────────┐
    │                        피아 식별 결과                             │
    └───────────────────────────┬─────────────────────────────────────┘
                                │
            ┌───────────────────┴───────────────────┐
            ▼                                       ▼
    ┌───────────────┐                      ┌───────────────┐
    │   👤 아군     │                      │   💀 적군     │
    └───────┬───────┘                      └───────┬───────┘
            │                                      │
            ▼                                      ▼
    "암구호를 말하라"                        "암구호를 말하라"
            │                                      │
      ┌─────┴─────┐                          ┌─────┴─────┐
      ▼           ▼                          ▼           ▼
   ✅ 정답     ❌ 오답                    ✅ 정답     ❌ 오답
      │           │                          │           │
      ▼           ▼                          ▼           ▼
   "통과"     📸 웹캠 캡처              📸 웹캠 캡처   🔫 헤드샷
   음성출력    + 웹 UI 알림             + 웹 UI 경고    발사!
    """)
    
    # ==================== 슬라이드 4: 프로젝트 소개 ====================
    add_section_slide(prs, "프로젝트 소개", "1️⃣")
    
    add_content_slide(prs, "🎯 프로젝트 목표", [
        '"로봇이 사람 얼굴을 실시간으로 추적하며 바라본다"',
        "",
        "💡 핵심 가치:",
        "   • 실시간성: 30fps 감지, <100ms 반응",
        "   • 안정성: EKF 기반 노이즈 제거",
        "   • 단순성: IK 없이 조인트 직접 제어",
        "",
        "🎬 [데모 영상/GIF 삽입 위치]"
    ], bullet=False)
    
    # ==================== 슬라이드 4-5: 문제 정의 ====================
    add_section_slide(prs, "문제 정의", "2️⃣")
    
    add_diagram_slide(prs, "❓ 얼굴 추적이 왜 어려운가?", """
    ┌──────────────┐     ┌──────────────┐     ┌──────────────┐
    │   얼굴 감지   │ ──▶ │  3D 위치 추정  │ ──▶ │   로봇 제어   │
    │  (2D 이미지)  │     │  (깊이 정보)   │     │ (6축 매니퓰레이터)│
    └──────────────┘     └──────────────┘     └──────────────┘
           ↓                    ↓                    ↓
       🔴 노이즈 多          🔴 좌표계 복잡         🔴 특이점 문제
       🔴 조명 변화          🔴 깊이 오차           🔴 IK 계산 복잡
       🔴 측면 얼굴          🔴 동기화 문제         🔴 반응 지연


    🎯 우리가 풀고자 한 문제:
       1. 빠른 감지: 실시간 30fps 이상
       2. 정확한 추정: 3D 위치 오차 최소화
       3. 부드러운 제어: 떨림 없는 로봇 동작
       4. 안전한 동작: 특이점/충돌 방지
    """)
    
    # ==================== 슬라이드 6-12: 개발 히스토리 ====================
    add_section_slide(prs, "개발 히스토리", "3️⃣")
    
    add_diagram_slide(prs, "📅 전체 타임라인", """
    Dec 8 (Day 1)  ──▶  Dec 9 (Day 2)  ──▶  Dec 10 (Day 3)  ──▶  Dec 12-13 (Day 4)
        │                   │                    │                      │
        ▼                   ▼                    ▼                      ▼
    ┌─────────┐      ┌───────────┐      ┌────────────────┐      ┌─────────────────┐
    │ 환경 구축 │      │ YOLO 학습  │      │ MediaPipe+MPC  │      │ YOLOv8+EKF+Joint│
    │캘리브레이션│      │ CNN 이론   │      │ 얼굴추적 v1    │      │ 최종 시스템      │
    └─────────┘      └───────────┘      └────────────────┘      └─────────────────┘
    
    
    🔧 기술 스택 발전:
    
       감지: Haar ─────▶ MediaPipe ─────▶ YOLOv8 + TensorRT
       추적: 없음 ─────▶ Moving Avg ────▶ EKF 9-state
       제어: 없음 ─────▶ Cartesian  ────▶ Joint-space
    """)
    
    # Day 1
    add_table_slide(prs, "📅 Day 1: 환경 구축 및 캘리브레이션 (Dec 8)", 
        ["작업", "상세", "결과"],
        [
            ["카메라 설정", "RealSense D435i 런치", "✅ 30Hz RGB+Depth"],
            ["Modbus 통신", "OnRobot RG2 그리퍼", "✅ 열기/닫기 제어"],
            ["캘리브레이션", "Eye-on-Hand", "✅ T_gripper2camera.npy"],
            ["URDF 작성", "로봇+그리퍼+카메라", "✅ TF 트리 완성"],
        ]
    )
    
    # Day 2
    add_content_slide(prs, "📅 Day 2: 딥러닝 기초 및 YOLO 학습 (Dec 9)", [
        "📚 이론 학습:",
        "   • CNN 구조: Conv → Pool → FC",
        "   • Object Detection: R-CNN → YOLO 발전 과정",
        "   • YOLO 원리: Single-shot detection",
        "",
        "💻 실습:",
        "   • YOLOv8 학습: Ultralytics 프레임워크",
        "   • 추론 테스트: 실시간 객체 검출",
        "",
        "💡 인사이트:",
        '   "얼굴 검출에는 얼굴 특화 모델이 필요하다"',
        "   → YOLOv8n-face 모델 발견!"
    ], bullet=False)
    
    # Day 3
    add_content_slide(prs, "📅 Day 3: 실시간 추적 시스템 v1 (Dec 10)", [
        "🔴 문제 1: Haar Cascade 정확도 60-70%",
        "   → 해결: MediaPipe 교체 → 90%+ 달성",
        "",
        "🔴 문제 2: Tracking Node 병목 (2-4Hz)",
        "   → 원인: 타이머 0.1초, TF2 블로킹",
        "   → 해결: 타이머 0.033초, 타임아웃 단축",
        "   → 결과: 2-4Hz → 30.3Hz (10배 향상!)",
        "",
        "🔴 문제 3: MPC 제어기 구현 시도",
        "   → N=10 예측, Q=100 추적, R=1 제어",
        "   → 하지만... 특이점/IK 문제 발생!"
    ], bullet=False)
    
    # Day 4
    add_content_slide(prs, "📅 Day 4: 최종 시스템 완성 (Dec 12-13)", [
        "⚡ Phase 1: YOLO + TensorRT",
        "   • YOLOv8n-face → TensorRT 엔진 변환",
        "   • 추론 시간: 7.6ms (FP16)",
        "",
        "⚡ Phase 2: MPC/Cartesian 포기 결정 🚨",
        "   • movel() 기반 → 특이점 문제 발생!",
        "   • IK 계산 지연 → 실시간성 저하",
        "",
        "⚡ Phase 3: Joint-space 직접 제어 (해결책)",
        '   • 핵심: "얼굴 추적 = 방향 추적 문제"',
        "   • J1(수평) + J4(수직) 직접 제어 → IK 불필요!",
        "",
        "⚡ Phase 4: 9-state EKF → 노이즈 80% 감소"
    ], bullet=False)
    
    # ==================== MPC → Joint 전환 슬라이드 ====================
    add_section_slide(prs, "MPC → Joint 전환", "🔄")
    
    add_diagram_slide(prs, "🚨 MPC 제어기의 문제점", """
    ┌─────────────────────────────────────────────────────────────────┐
    │                    MPC (Model Predictive Control)               │
    │                                                                 │
    │   목표 위치 (x,y,z) → IK 계산 → 관절각 → 로봇 이동              │
    └─────────────────────────────────────────────────────────────────┘
                                    │
                        ┌───────────┴───────────┐
                        ▼                       ▼
               ┌─────────────┐         ┌─────────────┐
               │ 🔴 특이점   │         │ 🔴 IK 지연  │
               │   문제      │         │    문제     │
               └─────────────┘         └─────────────┘
                      │                       │
                      ▼                       ▼
             팔이 완전히 펴지거나      실시간 30Hz에서
             접히면 IK 해 없음         IK 계산이 병목

    💡 Day 3 시도한 MPC 파라미터:
       • 예측 구간: N = 10 스텝
       • 추적 가중치: Q = 100 (위치 오차 패널티)
       • 제어 가중치: R = 1 (제어 입력 패널티)
       • 부드러움: S = 10 (가속도 변화 패널티)
       → 결과: 간헐적 특이점 진입, 불안정한 동작
    """)
    
    add_diagram_slide(prs, "💡 발상의 전환: 위치 → 방향", """
    ┌─────────────────────────────────────────────────────────────────┐
    │  기존 MPC 접근 (위치 추적)                                       │
    │                                                                 │
    │    "로봇 End-Effector를 얼굴 위치(x,y,z)로 이동시켜라"           │
    │    → 6축 IK 계산 필요 → 특이점 위험 → 복잡함                     │
    └─────────────────────────────────────────────────────────────────┘
                                    ▼
                              🔄 발상 전환!
                                    ▼
    ┌─────────────────────────────────────────────────────────────────┐
    │  새로운 Joint 접근 (방향 추적)                                   │
    │                                                                 │
    │    "로봇이 얼굴 방향만 바라보면 된다!"                            │
    │    → 2축만 제어 (J1: 수평, J4: 수직)                            │
    │    → IK 불필요 → 특이점 없음 → 단순함                           │
    └─────────────────────────────────────────────────────────────────┘

    🎯 핵심 인사이트:
       "로봇팔 끝에 레이저 포인터가 달려있다고 생각하면,
        레이저가 얼굴을 가리키기만 하면 됨!"
    """)
    
    add_table_slide(prs, "📊 MPC vs Joint-space 비교",
        ["항목", "MPC (Day 3)", "Joint-space (Day 4)"],
        [
            ["IK 계산", "필요 (6축)", "불필요"],
            ["특이점 위험", "있음 (치명적)", "없음"],
            ["제어 주기", "~10Hz", "50Hz"],
            ["구현 복잡도", "높음", "낮음"],
            ["코드 라인", "~300줄", "~100줄"],
            ["안정성", "불안정", "매우 안정"],
        ]
    )
    
    # ==================== 슬라이드 13-14: 시스템 아키텍처 ====================
    add_section_slide(prs, "시스템 아키텍처", "4️⃣")
    
    add_diagram_slide(prs, "🏗️ 전체 파이프라인", """
    ┌─────────────────────────────────────────────────────────────────────┐
    │                         RealSense D435i                              │
    │                      (RGB 640x480 + Depth)                           │
    └───────────────────────────────┬─────────────────────────────────────┘
                                    │ 30Hz
                                    ▼
    ┌─────────────────────────────────────────────────────────────────────┐
    │                     Face Detection Node                              │
    │     CLAHE 전처리 → YOLOv8-face → ROI Tracking → Confidence Filter   │
    │                        TensorRT FP16 (7.6ms)                         │
    └───────────────────────────────┬─────────────────────────────────────┘
                                    │ /face_detection/faces
                                    ▼
    ┌─────────────────────────────────────────────────────────────────────┐
    │                     Face Tracking Node                               │
    │      Depth Extraction → TF2 Transform → EKF 9-state → Marker        │
    │              camera_link → base_link 좌표 변환                        │
    └───────────────────────────────┬─────────────────────────────────────┘
                                    │ /face_tracking/marker_robot
                                    ▼
    ┌─────────────────────────────────────────────────────────────────────┐
    │                    Joint Tracking Node                               │
    │     Angle Calc → Dead Zone → Velocity Control → Jog Multi Axis      │
    │                  J1(수평) + J4(수직) 직접 제어                         │
    └───────────────────────────────┬─────────────────────────────────────┘
                                    ▼
    ┌─────────────────────────────────────────────────────────────────────┐
    │                       Doosan M0609                                   │
    └─────────────────────────────────────────────────────────────────────┘
    """)
    
    # ==================== 슬라이드 15-17: 기술 선택 비교 ====================
    add_section_slide(prs, "기술 선택 비교", "5️⃣")
    
    add_table_slide(prs, "🔍 얼굴 감지 기술 비교",
        ["기술", "속도", "정확도", "GPU", "선택"],
        [
            ["Haar Cascade", "30fps", "60%", "❌", "❌ 정확도 부족"],
            ["dlib HOG", "10fps", "80%", "❌", "❌ 느림"],
            ["MediaPipe", "60fps", "90%", "△", "△ Day3"],
            ["YOLOv8-face", "30fps", "95%", "✅", "✅ 최종 선택"],
        ]
    )
    
    add_table_slide(prs, "🔍 추적 필터 비교",
        ["필터", "예측", "노이즈 제거", "지연", "선택"],
        [
            ["Moving Average", "❌", "약함", "있음", "❌"],
            ["Low-pass Filter", "❌", "중간", "있음", "❌"],
            ["Kalman (6-state)", "속도만", "강함", "적음", "△"],
            ["EKF (9-state)", "가속도", "매우 강함", "최소", "✅"],
        ]
    )
    
    add_table_slide(prs, "🔍 로봇 제어 방식 비교",
        ["방식", "IK 필요", "특이점", "반응속도", "선택"],
        [
            ["Cartesian movel()", "✅", "위험", "느림", "❌"],
            ["MoveIt", "✅", "안전", "느림", "❌"],
            ["Joint-space Jog", "❌", "없음", "빠름", "✅"],
        ]
    )
    
    # ==================== 슬라이드 18-19: 핵심 알고리즘 ====================
    add_section_slide(prs, "핵심 알고리즘", "6️⃣")
    
    add_diagram_slide(prs, "🧮 Extended Kalman Filter (9-state)", """
    상태 벡터 (9차원):
        x = [x, y, z, vx, vy, vz, ax, ay, az]ᵀ
             ├─위치─┤ ├──속도──┤ ├──가속도─┤

    상태 전이 방정식 (등가속도 모델):
        p_k = p_{k-1} + v_{k-1}·dt + 0.5·a_{k-1}·dt²
        v_k = v_{k-1} + a_{k-1}·dt
        a_k = a_{k-1}  (가속도 일정 가정)

    상태 전이 행렬 F:
        F = | I   dt·I   0.5·dt²·I |
            | 0    I      dt·I      |
            | 0    0       I        |

    ✅ 효과:
       • 노이즈 80% 감소
       • 50ms 미래 위치 예측
       • 빠른 움직임도 부드럽게 추적
    """)
    
    add_diagram_slide(prs, "🎮 Joint-space 제어 알고리즘", """
    💡 핵심 인사이트:
       "얼굴 추적 = 방향 추적 문제"
       → 로봇이 특정 위치로 갈 필요 없음
       → 얼굴 방향만 바라보면 됨!

    🎯 제어 전략:
       J1: 수평 방향 (베이스 회전) - 메인
       J4: 수직 방향 (손목 피치) - 서브
       J2, J3, J5, J6: 고정 (팔 자세 유지)

    📐 제어 로직:
       error_horizontal = face_y  (로봇 Y축)
       j1_velocity = K_p × error_horizontal

       error_vertical = face_z    (로봇 Z축)
       j4_velocity = K_p × error_vertical

    ⚙️ 파라미터:
       K_p = 0.5, Dead Zone = 2°
       Max Vel: J1=30°/s, J4=40°/s
       Control Rate: 50Hz
    """)
    
    # ==================== 슬라이드: 결과 및 성과 ====================
    add_section_slide(prs, "결과 및 성과", "7️⃣")
    
    add_table_slide(prs, "📊 정량적 성과",
        ["지표", "목표", "달성", "달성률"],
        [
            ["감지 속도", "20fps", "30fps", "✅ 150%"],
            ["감지 정확도", "90%", "95%+", "✅ 달성"],
            ["추론 시간", "20ms", "7.6ms", "✅ 263%"],
            ["반응 시간", "200ms", "<100ms", "✅ 200%"],
            ["추적 거리", "1~2m", "0.3~3m", "✅ 150%"],
            ["노이즈 감소", "50%", "80%", "✅ 160%"],
        ]
    )
    
    add_content_slide(prs, "✅ 정성적 성과", [
        "조명 변화에도 안정적 감지",
        "빠른 머리 움직임도 부드럽게 추적",
        "특이점 문제 없이 연속 동작",
        "멀리서도(3m) 정확한 추적",
        "지터링/떨림 최소화",
        "",
        "📈 Day 1 → Day 4 발전:",
        "   • 감지: 60% → 95% (+35%)",
        "   • 추적: 2Hz → 30Hz (15배)",
        "   • 지연: 250ms → 100ms (2.5배)"
    ])
    
    # ==================== 미완성 모듈 & 향후 계획 ====================
    add_section_slide(prs, "미완성 모듈 & 향후 계획", "8️⃣")
    
    add_table_slide(prs, "📦 전체 시스템 완성도",
        ["모듈", "담당", "상태", "비고"],
        [
            ["헤드샷 트래킹", "태슬라", "✅ 완료", "YOLOv8+EKF+Joint"],
            ["SAM3 파지", "성우", "🔄 진행중", "권총 그립 인식"],
            ["피아 식별", "경훈", "🔄 진행중", "아군/적군 분류"],
            ["음성 암구호", "지원", "🔄 진행중", "STT/TTS 연동"],
            ["웹 UI", "TBD", "⏳ 대기", "알림/경고 인터페이스"],
            ["시스템 통합", "All", "⏳ 대기", "전체 시퀀스 연결"],
        ]
    )
    
    add_content_slide(prs, "🔧 미완성 모듈 상세", [
        "📌 SAM3 기반 파지 (sam3_grip_detection)",
        "   • Segment Anything Model 3 활용",
        "   • 권총 그립 포인트 자동 인식",
        "   • OnRobot RG2 그리퍼 연동 필요",
        "",
        "📌 피아 식별 모듈",
        "   • 아군/적군 이미지 학습 데이터 구축",
        "   • YOLOv8 Classification 또는 Face Recognition",
        "   • 헤드샷 트래킹과 연동 필요",
        "",
        "📌 음성 암구호 시스템",
        "   • STT: 음성 → 텍스트 (Whisper 등)",
        "   • TTS: 텍스트 → 음성 (gTTS 등)",
        '   • 암구호 판정: "오늘의 암구호는?"'
    ], bullet=False)
    
    add_diagram_slide(prs, "🔗 통합 시퀀스 (TODO)", """
    ┌───────────────────────────────────────────────────────────────────┐
    │                     미구현 (SAM3 팀 담당)                          │
    │  ┌─────────┐    ┌─────────────┐    ┌─────────────┐               │
    │  │ 권총 감지│ → │ SAM3 그립   │ → │ 파지 & 사로 │               │
    │  │         │    │ 포인트 인식 │    │ 이동        │               │
    │  └─────────┘    └─────────────┘    └─────────────┘               │
    └───────────────────────────────────────────────────────────────────┘
                                    ▼
    ┌───────────────────────────────────────────────────────────────────┐
    │                     ✅ 구현 완료 (태슬라)                          │
    │  ┌─────────────┐    ┌─────────────┐    ┌─────────────┐           │
    │  │ 얼굴 감지   │ → │ EKF 추적    │ → │ Joint 제어  │           │
    │  │ YOLOv8-face │    │ 9-state     │    │ J1+J4       │           │
    │  └─────────────┘    └─────────────┘    └─────────────┘           │
    └───────────────────────────────────────────────────────────────────┘
                                    ▼
    ┌───────────────────────────────────────────────────────────────────┐
    │                     미구현 (경훈 + 지원 담당)                      │
    │  ┌─────────────┐    ┌─────────────┐    ┌─────────────┐           │
    │  │ 피아 식별   │ → │ 암구호 판정 │ → │ 발사/경고   │           │
    │  │ 아군/적군   │    │ STT + TTS   │    │ 웹 UI       │           │
    │  └─────────────┘    └─────────────┘    └─────────────┘           │
    └───────────────────────────────────────────────────────────────────┘
    """)
    
    add_content_slide(prs, "🚀 향후 통합 계획", [
        "📅 Week 1: 개별 모듈 완성",
        "   • SAM3 파지: 권총 그립 인식 완성",
        "   • 피아 식별: 학습 데이터 & 모델",
        "   • 음성: STT/TTS 파이프라인",
        "",
        "📅 Week 2: 시스템 통합",
        "   • ROS2 Topic/Service 인터페이스 정의",
        "   • 모듈간 통신 테스트",
        "   • 전체 시퀀스 연결",
        "",
        "📅 Week 3: 테스트 & 데모",
        "   • 시나리오 테스트",
        "   • 예외 상황 처리",
        "   • 최종 데모 준비"
    ], bullet=False)
    
    # ==================== 슬라이드 24: 기술 스택 ====================
    add_content_slide(prs, "📚 기술 스택 요약", [
        "🔧 Hardware:",
        "   Doosan M0609 / RealSense D435i / NVIDIA RTX 4060",
        "",
        "💻 Software:",
        "   ROS2 Humble / Python 3.10 / TensorRT",
        "",
        "🧠 Algorithms:",
        "   YOLOv8n-face / 9-state EKF / Joint-space Control",
        "",
        "📦 패키지 구조:",
        "   face_tracking/ (1,837줄)",
        "   ├── detection/ (face_detection_node, yolo_detector)",
        "   ├── tracking/ (face_tracking_node, ekf_filter)",
        "   └── control/ (joint_tracking_node)"
    ], bullet=False)
    
    # ==================== 슬라이드 25: 감사 ====================
    add_title_slide(
        prs,
        "🙏 감사합니다",
        "Q&A\n\nGitHub: github.com/taesla/rokey_c_1_collabo2\nBranch: ros-face-tracking"
    )
    
    # 저장
    output_path = os.path.join(os.path.dirname(__file__), "FaceTracking_Presentation.pptx")
    prs.save(output_path)
    print(f"✅ PPT 생성 완료: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
